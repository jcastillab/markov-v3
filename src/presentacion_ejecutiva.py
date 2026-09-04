"""Generador de presentacion ejecutiva para el proyecto Markov Freedom."""

from __future__ import annotations

import json
import tempfile
from pathlib import Path

import matplotlib
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap
from pptx.util import Inches, Pt

try:
    from models.bayes import HierarchicalNB
    from models.supervised import build_supervised_dataset, feature_groups
except ModuleNotFoundError:
    from src.models.bayes import HierarchicalNB
    from src.models.supervised import build_supervised_dataset, feature_groups

matplotlib.use("Agg")

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "reports"
OUT.mkdir(parents=True, exist_ok=True)


def load_predictions():
    """Carga los mejores modelos y sus predicciones walk-forward."""
    eval_dir = ROOT / "outputs" / "evaluation"
    # El RF rolling usa los hiperparametros ganadores de la validacion fija.
    files = {
        "M3 BASE": (eval_dir / "predictions_e00_m3_base_rolling.csv", "proyectado"),
        "Mejor RF": (eval_dir / "predictions_modelo_seleccionado_rolling.csv", "proyectado"),
    }
    frames = {}
    for name, (path, pred_col) in files.items():
        if not path.exists():
            continue
        df = pd.read_csv(path)
        df["fecha_objetivo"] = pd.to_datetime(df["fecha_objetivo"])
        df["modelo_nombre"] = name
        df["scope"] = "ROLLING"
        pred_col_real = next(
            (column for column in (pred_col, "proyectado", "pred", "pred_bloque")
             if column in df.columns),
            None,
        )
        if pred_col_real is None:
            continue
        df = df.rename(columns={pred_col_real: "proyectado"})
        frames[name] = df
    bayes = build_bayes_rolling()
    if not bayes.empty:
        bayes["modelo_nombre"] = "Mejor Bayesiano"
        bayes["modelo"] = "Mejor Bayesiano"
        frames["Mejor Bayesiano"] = bayes
    return frames


def build_bayes_rolling():
    """Genera predicciones NB walk-forward usando solo origenes anteriores."""
    datasets = ROOT / "outputs" / "datasets"
    required = ["forecast_windows.parquet", "fact_bloque_dia.parquet",
                "transition_intervals_tradicional.parquet", "poda_features.parquet",
                "clima_features.parquet"]
    if not all((datasets / name).exists() for name in required):
        return pd.DataFrame()
    cfg_path = ROOT / "config" / "pipeline.yaml"
    from canonical import load_config
    cfg = load_config(cfg_path)
    windows = pd.read_parquet(datasets / required[0])
    fact = pd.read_parquet(datasets / required[1])
    intervals = pd.read_parquet(datasets / required[2])
    pruning = pd.read_parquet(datasets / required[3])
    climate = pd.read_parquet(datasets / required[4])
    frame = build_supervised_dataset(windows, fact, intervals, cfg, pruning, climate,
                                     include_incomplete=True)
    frame["fecha_origen"] = pd.to_datetime(frame["fecha_origen"])
    rows = []
    for origin in sorted(frame.fecha_origen.dropna().unique()):
        current = frame.fecha_origen.eq(origin)
        previous = frame.fecha_origen.lt(origin) & frame.target.notna()
        if previous.sum() == 0 or not current.any():
            continue
        model = HierarchicalNB(cfg["bayes"]["hierarchical_shrinkage"]).fit(frame.loc[previous])
        pred = model.predict(frame.loc[current])
        current_frame = frame.loc[current, ["finca", "bloque", "fecha_origen",
                                             "fecha_objetivo", "semana_proyeccion",
                                             "horizonte_dia", "target"]].copy()
        current_frame = current_frame.rename(columns={"target": "real"})
        current_frame["proyectado"] = pred
        current_frame["modelo_nombre"] = "NB jerarquico"
        current_frame["scope"] = "ROLLING"
        rows.append(current_frame)
    return pd.concat(rows, ignore_index=True) if rows else pd.DataFrame()


def build_weekly(frames):
    """Agrega predicciones a nivel finca-semana para todas las semanas."""
    weekly = []
    for name, df in frames.items():
        if "semana_proyeccion" not in df.columns:
            df["semana_proyeccion"] = (
                df["fecha_objetivo"].dt.isocalendar().year * 100
                + df["fecha_objetivo"].dt.isocalendar().week
            )
        agg = df.groupby(["finca", "semana_proyeccion"], as_index=False).agg(
            real=("real", lambda s: s.sum(min_count=1)),
            proyectado=("proyectado", lambda s: s.sum(min_count=1)),
        )
        agg["modelo"] = name
        weekly.append(agg)
    return pd.concat(weekly, ignore_index=True)


def champion_table(weekly):
    """Tabla ancha con Real, M3, RF, Bayes por finca y semana."""
    wide = weekly.pivot_table(
        index=["finca", "semana_proyeccion"],
        columns="modelo",
        values="proyectado",
        aggfunc="first",
    ).reset_index()
    real = (
        weekly.groupby(["finca", "semana_proyeccion"], as_index=False)["real"]
        .first()
        .rename(columns={"real": "Conteo real"})
    )
    wide = real.merge(wide, on=["finca", "semana_proyeccion"], how="left")
    rename = {
        "M3 BASE": "M3 BASE",
        "Mejor RF": "Mejor RF",
        "Mejor Bayesiano": "Mejor Bayesiano",
        "M3 Dirichlet": "M3 Dirichlet",
    }
    wide = wide.rename(columns=rename)
    wide = wide[["finca", "semana_proyeccion", "Conteo real", "M3 BASE", "Mejor RF", "Mejor Bayesiano"]]
    return wide


def hit_table(weekly):
    """Resume todas las semanas observadas con la escala operativa."""
    df = weekly.copy()
    df["error_relativo"] = (df["proyectado"] - df["real"]).abs() / df["real"].abs()
    df.loc[df["real"].eq(0), "error_relativo"] = np.nan
    df["estado_acierto"] = np.select(
        [df.error_relativo.le(0.07), df.error_relativo.le(0.10)],
        ["ACIERTO", "CERCA"], default="NO ACIERTO")
    df.loc[df.error_relativo.isna(), "estado_acierto"] = "SIN REAL"
    observed = df[df.estado_acierto.ne("SIN REAL")]
    summary = observed.groupby(["finca", "modelo"], as_index=False).agg(
        semanas_acertadas=("estado_acierto", lambda s: s.eq("ACIERTO").sum()),
        semanas_cerca=("estado_acierto", lambda s: s.eq("CERCA").sum()),
        semanas_no_acierto=("estado_acierto", lambda s: s.eq("NO ACIERTO").sum()),
        semanas_evaluables=("estado_acierto", "size"),
    )
    summary["porcentaje_acierto"] = summary.semanas_acertadas / summary.semanas_evaluables
    return summary


def plot_historico(weekly):
    """Grafico de lineas historicas por finca."""
    fig, axes = plt.subplots(1, 3, figsize=(16, 5), sharey=True)
    colors = {"Conteo real": "#2c3e50", "M3 BASE": "#c0392b", "Mejor RF": "#27ae60", "Mejor Bayesiano": "#2980b9"}
    for ax, finca in zip(axes, sorted(weekly["finca"].unique())):
        sub = weekly[weekly.finca == finca].copy()
        sub["semana_label"] = sub["semana_proyeccion"].astype(str)
        for modelo in ["Conteo real", "M3 BASE", "Mejor RF", "Mejor Bayesiano"]:
            model_name = (
                modelo
                if modelo == "Conteo real"
                else {"M3 BASE": "M3 BASE", "Mejor RF": "Mejor RF", "Mejor Bayesiano": "Mejor Bayesiano"}[modelo]
            )
            if modelo == "Conteo real":
                data = sub.drop_duplicates("semana_proyeccion")[["semana_label", "real"]].rename(columns={"real": "value"})
            else:
                data = sub[sub.modelo == model_name][["semana_label", "proyectado"]].rename(columns={"proyectado": "value"})
            if not data.empty:
                ax.plot(data["semana_label"], data["value"], label=modelo,
                        color=colors[modelo], marker="o", markersize=3, linewidth=2)
        ax.set_title(finca)
        ax.tick_params(axis="x", rotation=45)
        ax.set_ylabel("Tallos")
        ax.legend(fontsize=7)
        ax.grid(True, linestyle="--", alpha=0.4)
    plt.tight_layout()
    path = OUT / "historico_modelos.png"
    plt.savefig(path, dpi=150)
    plt.close()
    return path


def plot_model_comparison(ranking):
    """Grafico ejecutivo de WAPE y acierto global."""
    selected = ranking.head(5)
    if selected.empty:
        return None
    fig, ax = plt.subplots(figsize=(10, 4.5))
    names = selected.experiment_id.str.replace("_", " ")
    bars = ax.barh(names, selected.wape * 100,
                   color=["#c0392b", "#27ae60", "#2980b9", "#8e44ad"][:len(selected)])
    ax.set_xlabel("WAPE (%) | menor es mejor")
    ax.set_title("Desempeño causal comparable")
    ax.grid(axis="x", linestyle="--", alpha=.35)
    for bar, value in zip(bars, selected.wape * 100):
        ax.text(value + .5, bar.get_y() + bar.get_height() / 2,
                f"{value:.1f}%", va="center", fontsize=10)
    fig.tight_layout()
    path = OUT / "comparacion_modelos.png"
    fig.savefig(path, dpi=160)
    plt.close(fig)
    return path


def load_feature_inventory():
    """Resume las familias de variables disponibles en el dataset supervisado."""
    path = ROOT / "outputs" / "datasets" / "dataset_supervisado_diario.parquet"
    frame = pd.read_parquet(path)
    numeric = set(frame.select_dtypes(include=[np.number]).columns)
    groups = {name: list(dict.fromkeys(col for col in cols if col in numeric and col != "target"))
              for name, cols in feature_groups(frame).items()}
    feno = set(groups["FENO"])
    return frame, groups, pd.DataFrame([
        {"Familia": "Fenología, M3, escala e historia", "Variables": len(feno), "Detalle": "conteos t0, proporciones, M3, camas y cortes rezagados"},
        {"Familia": "Clima causal", "Variables": len(set(groups["FENO_CLIMA"]) - feno), "Detalle": "VPD, GDD, DLI, temperatura, lluvia y ET0 en ventanas 1-56d"},
        {"Familia": "Podas y alineamiento", "Variables": len(set(groups["FENO_PODA"]) - feno), "Detalle": "lags, acumulados, kernels y días desde operación"},
        {"Familia": "FENO_CLIMA total", "Variables": len(groups["FENO_CLIMA"]), "Detalle": "familia usada por el RF operativo exploratorio"},
    ])


def plot_feature_composition(inventory):
    fig, ax = plt.subplots(figsize=(10, 4.8))
    base = inventory.iloc[0]
    climate = inventory.iloc[1]
    pruning = inventory.iloc[2]
    values = [base.Variables, climate.Variables, pruning.Variables]
    labels = [f"Base FENO\n({int(base.Variables)})",
              f"Clima adicional\n({int(climate.Variables)})",
              f"Poda adicional\n({int(pruning.Variables)})"]
    bars = ax.bar(labels, values, color=["#2c3e50", "#2980b9", "#d99b23"])
    ax.set_ylabel("Número de variables numéricas")
    ax.set_title("Cobertura de señales por familia de features")
    ax.grid(axis="y", linestyle="--", alpha=.35)
    for bar, value in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, value + 2, str(value), ha="center", fontweight="bold")
    ax.text(.5, -0.23, f"FENO_CLIMA = {int(base.Variables + climate.Variables)} variables", transform=ax.transAxes,
            ha="center", fontsize=10, color="#45505a")
    fig.tight_layout()
    path = OUT / "inventario_features.png"
    fig.savefig(path, dpi=160)
    plt.close(fig)
    return path


def plot_correlation_heatmap(frame):
    candidates = ["RC_t0", "SS_t0", "AP_t0", "p_RC", "p_SS", "p_AP", "log1p_RC",
                  "log1p_SS", "log1p_AP", "M3_pred_bloque", "corte_sum_7d", "corte_mean_7d"]
    cols = [col for col in candidates if col in frame.columns]
    corr_raw = frame[cols].corr().reindex(index=cols, columns=cols)
    corr_array = np.array(corr_raw, dtype=float, copy=True)
    np.fill_diagonal(corr_array, 1.0)
    corr = pd.DataFrame(corr_array, index=cols, columns=cols)
    fig, ax = plt.subplots(figsize=(10, 7.5))
    image = ax.imshow(corr, vmin=-1, vmax=1, cmap="RdBu")
    ax.set_xticks(range(len(cols)), cols, rotation=45, ha="right", fontsize=8)
    ax.set_yticks(range(len(cols)), cols, fontsize=8)
    for row in range(len(cols)):
        for col in range(len(cols)):
            value = corr.iloc[row, col]
            ax.text(col, row, f"{value:.2f}", ha="center", va="center", fontsize=6,
                    color="white" if abs(value) > .55 else "black")
    fig.colorbar(image, ax=ax, label="Correlación Pearson")
    ax.set_title("Correlaciones entre señales estructurales del RF")
    fig.tight_layout()
    path = OUT / "correlaciones_features.png"
    fig.savefig(path, dpi=160)
    plt.close(fig)
    return path


def color_for_hit(error):
    """Colores de la escala operativa: <=7%, <=10%, >10%."""
    if error <= 0.07:
        return "2e8b57"
    if error <= 0.10:
        return "d99b23"
    return "c94c4c"


def set_cell_text(cell, text, bold=False, size=10):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(16)
    return slide


def add_explanation_slide(prs, title, subtitle, equation, left_title, left_points,
                          right_title, right_points, accent):
    """Slide visual con mecanismo y evidencia no redundante."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(.55), Inches(.35), Inches(12.2), Inches(.55))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(26)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    sub = slide.shapes.add_textbox(Inches(.58), Inches(.95), Inches(12), Inches(.35))
    sub.text_frame.text = subtitle
    sub.text_frame.paragraphs[0].font.size = Pt(12)
    sub.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x75, 0x7D, 0x89)
    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.6), Inches(1.55), Inches(5.1), Inches(4.95))
    left.fill.solid(); left.fill.fore_color.rgb = RGBColor.from_string(accent); left.line.fill.background()
    tf = left.text_frame; tf.clear(); tf.margin_left = Inches(.28); tf.margin_right = Inches(.25)
    p = tf.paragraphs[0]; p.text = equation; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 255)
    p = tf.add_paragraph(); p.text = left_title; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 255); p.space_before = Pt(16)
    for point in left_points:
        p = tf.add_paragraph(); p.text = point; p.level = 0; p.font.size = Pt(14); p.font.color.rgb = RGBColor(255, 255, 255); p.space_before = Pt(12)
    right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.05), Inches(1.55), Inches(6.65), Inches(4.95))
    right.fill.solid(); right.fill.fore_color.rgb = RGBColor(0xF3, 0xF6, 0xF8); right.line.color.rgb = RGBColor(0xD8, 0xE1, 0xE5)
    rtf = right.text_frame; rtf.clear(); rtf.margin_left = Inches(.3)
    p = rtf.paragraphs[0]; p.text = right_title; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    for point in right_points:
        p = rtf.add_paragraph(); p.text = point; p.font.size = Pt(15); p.font.color.rgb = RGBColor(0x45, 0x4D, 0x5A); p.space_before = Pt(15)
    return slide


def add_process_slide(prs):
    """Mapa visual del flujo completo de modelado."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(.55), Inches(.35), Inches(12), Inches(.6))
    title.text_frame.text = "De datos crudos a decisión operativa"
    title.text_frame.paragraphs[0].font.size = Pt(26); title.text_frame.paragraphs[0].font.bold = True
    stages = [("01", "Datos", "Conteos,\nfenología, podas, clima", "2C3E50"),
              ("02", "M3", "Estados y\ntransiciones", "C0392B"),
              ("03", "Riesgo", "Duración y\nhazard", "8E44AD"),
              ("04", "RF", "Corrección\nno lineal", "27AE60"),
              ("05", "Bayes", "Pooling e\nincertidumbre", "2980B9"),
              ("06", "Acción", "Mano de obra,\npacking y ventas", "D99B23")]
    for i, (num, name, text, color) in enumerate(stages):
        x = .45 + i * 2.12
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.1), Inches(1.8), Inches(2.25))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor.from_string(color); box.line.fill.background()
        tf = box.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)
        p = tf.add_paragraph(); p.text = name; p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)
        p = tf.add_paragraph(); p.text = text; p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(11); p.font.color.rgb = RGBColor(255,255,255)
        if i < len(stages) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.82), Inches(3.0), Inches(.3), Inches(.35))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = RGBColor(0xB8,0xC2,0xC8); arrow.line.fill.background()
    note = slide.shapes.add_textbox(Inches(.7), Inches(5.2), Inches(12), Inches(1.1))
    note.text_frame.text = "Regla transversal: toda predicción respeta causalidad temporal y se compara en la misma población cuando se decide el champion."
    note.text_frame.paragraphs[0].font.size = Pt(17); note.text_frame.paragraphs[0].font.bold = True; note.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x2C,0x3E,0x50)
    return slide


def add_table_slide(prs, title, df, fmt_conditional=False, value_cols=None, total_row=False):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    rows, cols = len(df) + 1 + int(total_row), len(df.columns)
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(12.0)
    height = Inches(6.0)
    height = min(Inches(6.0), Inches(0.35 * rows))
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    # Encabezado
    for j, col in enumerate(df.columns):
        set_cell_text(table.cell(0, j), col, bold=True, size=10)
        table.cell(0, j).fill.solid()
        table.cell(0, j).fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
        table.cell(0, j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Datos
    for i, (_, row) in enumerate(df.iterrows()):
        for j, col in enumerate(df.columns):
            val = row[col]
            if isinstance(val, float) and pd.isna(val):
                display = "—"
            elif isinstance(val, float):
                if fmt_conditional and col in (value_cols or []):
                    display = f"{int(round(val)):,.0f}"
                else:
                    display = f"{val:.1%}" if "porcentaje" in str(col).lower() else f"{val:,.0f}"
            else:
                display = str(val)
            set_cell_text(table.cell(i + 1, j), display, size=8)
            if fmt_conditional and col in (value_cols or []):
                real = row["Conteo real"]
                if pd.notna(real) and real != 0 and isinstance(val, float) and pd.notna(val):
                    error = abs(val - real) / abs(real)
                    rgb = color_for_hit(error)
                    table.cell(i + 1, j).fill.solid()
                    table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(int(rgb[:2], 16), int(rgb[2:4], 16), int(rgb[4:], 16))
                    table.cell(i + 1, j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return slide


def add_image_slide(prs, title, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.3), width=Inches(12.0))
    return slide


def main():
    evaluation = ROOT / "outputs" / "evaluation"
    ranking = pd.read_csv(evaluation / "ranking_final.csv")
    comparison = pd.read_csv(evaluation / "metrics_comparacion_final.csv")
    selected = json.loads((evaluation / "selected_model_manifest.json").read_text(encoding="utf-8"))

    def metric(experiment_id, prefer_split="VALIDATION"):
        rows = comparison[comparison.experiment_id.eq(experiment_id)]
        preferred = rows[rows.split.eq(prefer_split)]
        return (preferred if not preferred.empty else rows).sort_values("wape").iloc[0]

    m3 = metric("E00_M3_BASE")
    p32 = comparison[(~comparison.causal.astype(str).str.lower().eq("true"))].sort_values("wape").iloc[0]
    climate = comparison[comparison.experiment_id.str.startswith("M3_CLIMA", na=False)].sort_values("wape").iloc[0]
    bayes = metric("NB_JERARQUICO")
    selected_validation = selected["selection"]
    ranking_n = int(ranking.n.min())
    champion = ranking.iloc[0]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Portada
    add_title_slide(
        prs,
        "Pronóstico de Corte Comercial de Rosas Freedom",
        "Línea de trabajo Markov M3 → Riesgo → Random Forest → Bayesiano\n" "Resultados, aprendizajes y modelo champion provisional",
    )

    # Agenda
    add_bullet_slide(
        prs,
        "Agenda",
        [
            "1. Contexto y reglas de oro del proyecto",
            "2. Modelo mecanicista: Markov M3",
            "3. Modelo de riesgo: Semi Markov y P32",
            "4. Random Forest: el champion provisional",
            "5. Modelos bayesianos",
            "6. Comparación final y tabla operativa",
            "7. Histórico de pronósticos por finca",
        ],
    )

    # Contexto
    add_bullet_slide(
        prs,
        "Contexto y reglas de oro",
        [
            "Fincas: ALMER, LA PRADERA, SANTA HELENA; variedad Freedom.",
            "Métrica principal: WAPE (no MAPE).",
            "Validación temporal y causal: información máxima <= t0.",
            "M3 es baseline obligatorio; ningún challenger gana por complejidad.",
            "Se excluyen modelos retrospectivos del ranking causal.",
        ],
    )
    add_process_slide(prs)

    # M3
    add_bullet_slide(
        prs,
        "1. Modelo mecanicista: Markov M3",
        [
            "Cadena de estados fenológicos: RC → SS → AP → PC.",
            "Matrices de transición por finca y vigencia (Abril / Julio).",
            "Simulación diaria del vector de estado a partir del conteo t0.",
            "Extrapolación muestra → bloque con factor de camas.",
            f"Resultado baseline: WAPE {m3.wape:.2%} sobre {int(m3.n)} días de validación causal.",
        ],
    )
    add_explanation_slide(prs, "M3: una fotografía fenológica convertida en trayectoria",
                          "Modelo de referencia mecanicista y baseline obligatorio",
                          "x(t+1) = Q x(t) + ingreso_RC",
                          "Qué usa",
                          ["3 conteos fenológicos: RC, SS y AP en t0.",
                           "Matriz por finca y vigencia Abril/Julio.",
                           "Exposiciones por grupos de duración: no requiere pares diarios."],
                          "Qué permite auditar",
                          ["Flujo RC → SS → AP → PC y pérdidas por estado.",
                           "AP→PC La Pradera: 47,78% Abril; 32,12% Julio.",
                            f"WAPE causal: {m3.wape:.2%}; sesgo: {m3.bias_pct:.2%}."], "C0392B")

    # Riesgo / P32
    add_bullet_slide(
        prs,
        "2. Modelo de riesgo: Semi Markov con P32",
        [
            "Aporta duración dentro del estado y edad latente del stock.",
            "P32 normalizado: 1.297 observaciones, 514 intervalos válidos.",
            "Se corrigió concentración artificial de edad inicial en cero.",
            f"Mejor resultado retrospectivo: {p32.experiment_id}, WAPE {p32.wape:.2%}.",
            "No es causal para ventanas históricas: P32 se levantó después.",
        ],
    )
    add_explanation_slide(prs, "Riesgo: no solo importa el estado, también la edad",
                          "La duración modifica la probabilidad de avanzar o cortar",
                          "P(siguiente | estado, edad)",
                          "Datos y parámetros",
                          ["1.297 observaciones P32; 514 intervalos válidos.",
                           "Variables: estado macro, edad, evento y exposición.",
                           "Hazards RC/SS/AP con suavizado hacia M3."],
                          "Lectura correcta del resultado",
                          [f"Mejor P32 retrospectivo: {p32.experiment_id}, WAPE {p32.wape:.2%}.",
                           "No se mezcla con el ranking: P32 fue capturado después de las ventanas.",
                           "Valor actual: aprender duración; no operar el resultado retrospectivo."], "8E44AD")

    # Podas y clima
    add_bullet_slide(
        prs,
        "Challengers exploratorios: podas y clima",
        [
            "Podas: features CORTE/ALINEAMIENTO con lags de 42-84 días.",
            "Podas no mejoraron WAPE; quedan documentadas como challenger.",
            "Clima LA PRADERA: VPD, GDD, DLI estimado, lluvia, ET0.",
            f"Mejor challenger M3-clima: WAPE {climate.wape:.2%} en su población local.",
            "Cobertura limitada a una finca y sin microclima de invernadero.",
        ],
    )
    feature_frame, feature_groups_map, inventory = load_feature_inventory()
    add_table_slide(prs, "Inventario de señales del modelo supervisado", inventory)
    add_image_slide(prs, "Familias de variables: base, clima y poda", plot_feature_composition(inventory))
    add_image_slide(prs, "Colinealidad en señales estructurales", plot_correlation_heatmap(feature_frame))

    # RF
    add_bullet_slide(
        prs,
        "3. Random Forest: champion provisional",
        [
            f"El modelo seleccionado es {selected['model']}.",
            f"Usa {selected['features']} y la familia {selected['family']}.",
            f"Selección temporal: WAPE diario {selected_validation['daily_wape']:.2%} y semanal {selected_validation['weekly_wape']:.2%}.",
            "La selección compuesta prioriza la escala semanal y se verifica en rolling-origin.",
            "Requiere validación congelada adicional antes de promoción operacional.",
        ],
    )
    add_explanation_slide(prs, "Random Forest: aprende la relación no lineal con el corte",
                           "Combina señales fenológicas, operativas, climáticas, historial y escala",
                           "corte = f(fenología, poda, clima, historial, exposición)",
                           "Variables usadas por el modelo seleccionado",
                           [f"Grupo seleccionado: {selected['features']}.",
                            "Base: conteos t0, proporciones, M3, camas e historial de cortes.",
                            "Señales adicionales: podas/alineamiento rezagados y clima causal."],
                           "Control y resultado",
                           ["Se detectó VIF alto en agregados de corte, proporciones y logs; RF tolera redundancia, pero reparte importancia.",
                            f"Selección temporal: WAPE diario {selected_validation['daily_wape']:.2%}; semanal {selected_validation['weekly_wape']:.2%}.",
                            "Validación temporal fija y rolling; falta tercer periodo congelado."], "27AE60")

    # Bayesiano
    add_bullet_slide(
        prs,
        "4. Modelos bayesianos",
        [
            "Dirichlet-Multinomial: posterior de matrices de transición centrado en M3.",
            "NB jerárquico: Gamma-Poisson con pooling por finca, bloque y horizonte.",
            f"NB jerárquico: WAPE {bayes.wape:.2%} en validación fija causal.",
            "Cobertura de intervalos inferior a la nominal: no usable para incertidumbre aún.",
            "Dirichlet-Multinomial no superó el baseline M3.",
        ],
    )
    add_explanation_slide(prs, "Bayes: regularización e incertidumbre explícita",
                          "Pooling entre finca, bloque y horizonte para evitar estimaciones frágiles",
                          "Y ~ NegBin(μ, α)",
                          "Estructura del modelo",
                          ["Nivel global → finca → bloque → horizonte.",
                           "Variables de agrupación: finca, bloque y día de horizonte.",
                           "Posterior Dirichlet adicional para matrices M3."],
                          "Resultado y limitación",
                           [f"NB jerárquico: WAPE {bayes.wape:.2%}, mejor Bayesiano causal.",
                            f"Cobertura: {bayes.coverage_interval_80:.1%} al 80% y {bayes.coverage_interval_95:.1%} al 95%; está descalibrada.",
                           "Aporta pooling para bloques escasos, no intervalos operativos todavía."], "2980B9")

    # Ranking final
    ranking_display = ranking.head(5)[["experiment_id", "wape", "acierto_global", "decision"]].copy()
    ranking_display["wape"] = ranking_display["wape"].apply(lambda x: f"{x:.2%}")
    ranking_display["acierto_global"] = ranking_display["acierto_global"].apply(lambda x: f"{x:.2%}")
    ranking_display.columns = ["Modelo", "WAPE", "Acierto", "Decisión"]
    add_table_slide(prs, f"5. Ranking rolling causal ({ranking_n} observaciones comunes)", ranking_display)
    comparison_path = plot_model_comparison(ranking)
    if comparison_path:
        add_image_slide(prs, "Lectura ejecutiva del desempeño", comparison_path)

    # Cargar predicciones
    frames = load_predictions()
    if frames:
        weekly = build_weekly(frames)
        wide = champion_table(weekly)
        # Mostrar todas las semanas disponibles en el backtest walk-forward.
        wide = wide.dropna(subset=["M3 BASE", "Mejor RF", "Mejor Bayesiano"], how="all")
        wide["semana_proyeccion"] = wide["semana_proyeccion"].astype(int)
        sample = wide.sort_values(["finca", "semana_proyeccion"])
        sample_display = sample.copy()
        sample_display.columns = ["Finca", "Semana", "Conteo real", "M3 BASE", "Mejor RF", "Mejor Bayesiano"]
        # Una tabla por finca evita comprimir 40+ semanas en una sola diapositiva.
        for finca, farm_table in sample_display.groupby("Finca", sort=True):
            add_table_slide(
                prs,
                f"6. Pronóstico histórico rolling: {finca}",
                farm_table,
                fmt_conditional=True,
                value_cols=["M3 BASE", "Mejor RF", "Mejor Bayesiano"],
            )

        hits = hit_table(weekly)
        hits_display = hits.copy()
        hits_display["porcentaje_acierto"] = hits_display["porcentaje_acierto"].apply(lambda x: f"{x:.1%}")
        hits_display.columns = ["Finca", "Modelo", "Acierto <=7%", "Cerca 7-10%", "No acierto >10%", "Semanas evaluables", "% acierto"]
        add_table_slide(prs, "7. Semanas acertadas por modelo y finca (escala operativa)", hits_display)

        img_path = plot_historico(weekly)
        add_image_slide(prs, "8. Histórico comparable: real vs mejores modelos", img_path)
    else:
        add_bullet_slide(prs, "6-8. Predicciones", ["No se encontraron archivos de predicciones en outputs/evaluation."])

    # Cierre
    add_bullet_slide(
        prs,
        "Conclusiones y próximos pasos",
        [
            f"El champion provisional rolling es {champion.experiment_id}, con WAPE {champion.wape:.2%}.",
            "M3 permanece como baseline mecanicista obligatorio y referencia.",
            "Bayesiano mejora el baseline pero aún no entrega incertidumbre calibrada.",
            "P32 y clima son prometedores pero requieren datos contemporáneos.",
            "Antes de producción: backtest congelado, validación por finca/bloque y trazabilidad.",
        ],
    )

    pptx_path = OUT / "presentacion_ejecutiva_markov_v3.pptx"
    try:
        prs.save(pptx_path)
    except PermissionError:
        pptx_path = OUT / "presentacion_ejecutiva_markov_v3_actualizada.pptx"
        prs.save(pptx_path)
    print(f"Presentacion guardada en: {pptx_path}")


if __name__ == "__main__":
    main()
